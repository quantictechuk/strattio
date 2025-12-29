"""Pitch Deck Generator routes - Auto-generate pitch decks from business plans"""

from fastapi import APIRouter, HTTPException, Depends, Response
from pydantic import BaseModel
from typing import Optional, Dict, List
from datetime import datetime
import logging
import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from utils.serializers import serialize_doc, to_object_id
from utils.dependencies import get_db
from utils.admin import get_current_user_id
from emergentintegrations.llm.chat import LlmChat, UserMessage

router = APIRouter()
logger = logging.getLogger(__name__)

class BrandingConfig(BaseModel):
    logo_url: Optional[str] = None
    primary_color: Optional[str] = "#001639"
    secondary_color: Optional[str] = "#3B82F6"
    font_family: Optional[str] = "Arial"

@router.post("/plans/{plan_id}/pitch-deck/generate")
async def generate_pitch_deck(
    plan_id: str,
    branding: Optional[BrandingConfig] = None,
    user_id: str = Depends(get_current_user_id),
    db = Depends(get_db)
):
    """Generate a pitch deck from a business plan"""
    
    # Get plan
    plan = await db.plans.find_one({
        "_id": to_object_id(plan_id),
        "user_id": user_id
    })
    
    if not plan:
        raise HTTPException(status_code=404, detail="Plan not found")
    
    # Get plan data
    sections = await db.sections.find({"plan_id": plan_id}).sort("order_index", 1).to_list(None)
    financial_model = await db.financial_models.find_one({"plan_id": plan_id})
    intake_data = plan.get("intake_data", {})
    
    # Extract key information
    business_name = intake_data.get("business_name", "Business")
    industry = intake_data.get("industry", "General")
    description = intake_data.get("business_description", "")
    value_prop = intake_data.get("unique_value_proposition", "")
    target_customers = intake_data.get("target_customers", "")
    revenue_model = intake_data.get("revenue_model", [])
    
    # Get section content
    exec_summary = next((s for s in sections if s.get("section_type") == "executive_summary"), None)
    market_analysis = next((s for s in sections if s.get("section_type") == "market_analysis"), None)
    team_section = next((s for s in sections if s.get("section_type") == "team"), None)
    
    # Generate slides using AI
    slides_data = await _generate_slides_content(
        plan, sections, financial_model, intake_data
    )
    
    # Create PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Set branding colors
    primary_color = RGBColor(0, 22, 57) if not branding or not branding.primary_color else _hex_to_rgb(branding.primary_color)
    secondary_color = RGBColor(59, 130, 246) if not branding or not branding.secondary_color else _hex_to_rgb(branding.secondary_color)
    accent_color = RGBColor(39, 172, 133)  # Green accent
    background_color = RGBColor(248, 250, 252)  # Light gray background
    text_color = RGBColor(71, 85, 105)  # Slate gray text
    
    # Create slides with professional design
    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add background rectangle for visual appeal
        if idx == 0:  # Title slide gets special treatment
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(10), Inches(2.5)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = primary_color
            bg_shape.line.fill.background()
        else:
            # Subtle top accent bar for other slides
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(10), Inches(0.3)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = secondary_color
            accent_bar.line.fill.background()
        
        # Add title with better styling
        if slide_data.get("title"):
            title_y = Inches(0.6) if idx == 0 else Inches(0.5)
            title_height = Inches(1.2) if idx == 0 else Inches(0.8)
            title_box = slide.shapes.add_textbox(Inches(0.75), title_y, Inches(8.5), title_height)
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44) if idx == 0 else Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255) if idx == 0 else primary_color
            title_para.alignment = PP_ALIGN.LEFT
            title_para.space_after = Pt(0)
        
        # Add content with better formatting
        if slide_data.get("content"):
            content_y = Inches(2.8) if idx == 0 else Inches(1.6)
            content_height = Inches(4.5) if idx == 0 else Inches(5.5)
            content_box = slide.shapes.add_textbox(Inches(0.75), content_y, Inches(8.5), content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.margin_left = Inches(0)
            content_frame.margin_right = Inches(0)
            content_frame.margin_top = Inches(0)
            content_frame.margin_bottom = Inches(0)
            
            # Parse content into paragraphs
            content_lines = slide_data["content"].split('\n')
            for line_idx, line in enumerate(content_lines):
                if line.strip():
                    if line_idx > 0:
                        p = content_frame.add_paragraph()
                    else:
                        p = content_frame.paragraphs[0]
                    
                    # Handle bullet points
                    if line.strip().startswith('•') or line.strip().startswith('-'):
                        p.text = line.strip().lstrip('•-').strip()
                        p.level = 0
                        p.font.size = Pt(20) if idx == 0 else Pt(18)
                        p.font.color.rgb = RGBColor(255, 255, 255) if idx == 0 else text_color
                        p.space_after = Pt(14)
                        p.space_before = Pt(0)
                    else:
                        p.text = line.strip()
                        p.level = 0
                        p.font.size = Pt(20) if idx == 0 else Pt(18)
                        p.font.color.rgb = RGBColor(255, 255, 255) if idx == 0 else text_color
                        p.space_after = Pt(14)
                        p.space_before = Pt(0)
        
        # Add slide number (except on title slide)
        if idx > 0:
            slide_num_shape = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
            slide_num_frame = slide_num_shape.text_frame
            slide_num_frame.text = f"{idx}"
            slide_num_para = slide_num_frame.paragraphs[0]
            slide_num_para.font.size = Pt(14)
            slide_num_para.font.color.rgb = RGBColor(148, 163, 184)
            slide_num_para.alignment = PP_ALIGN.RIGHT
    
    # Save to bytes
    deck_bytes = io.BytesIO()
    prs.save(deck_bytes)
    deck_bytes.seek(0)
    
    # Store deck metadata
    deck_doc = {
        "plan_id": plan_id,
        "user_id": user_id,
        "slides": slides_data,
        "branding": branding.dict() if branding else {},
        "created_at": datetime.utcnow(),
        "updated_at": datetime.utcnow()
    }
    
    await db.pitch_decks.insert_one(deck_doc)
    
    logger.info(f"Pitch deck generated for plan {plan_id}")
    
    return {
        "message": "Pitch deck generated successfully",
        "deck_id": str(deck_doc["_id"]),
        "slide_count": len(slides_data)
    }

@router.get("/plans/{plan_id}/pitch-deck/download")
async def download_pitch_deck(
    plan_id: str,
    format: str = "pptx",
    user_id: str = Depends(get_current_user_id),
    db = Depends(get_db)
):
    """Download the generated pitch deck"""
    
    # Get plan
    plan = await db.plans.find_one({
        "_id": to_object_id(plan_id),
        "user_id": user_id
    })
    
    if not plan:
        raise HTTPException(status_code=404, detail="Plan not found")
    
    # Get latest deck
    deck = await db.pitch_decks.find_one({
        "plan_id": plan_id,
        "user_id": user_id
    }, sort=[("created_at", -1)])
    
    if not deck:
        raise HTTPException(status_code=404, detail="Pitch deck not found. Please generate one first.")
    
    # Regenerate PPTX (for now, we regenerate each time)
    # In production, you'd store the file and serve it
    sections = await db.sections.find({"plan_id": plan_id}).sort("order_index", 1).to_list(None)
    financial_model = await db.financial_models.find_one({"plan_id": plan_id})
    intake_data = plan.get("intake_data", {})
    
    slides_data = await _generate_slides_content(plan, sections, financial_model, intake_data)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    primary_color = RGBColor(0, 22, 57)
    secondary_color = RGBColor(59, 130, 246)
    
    # Use the same improved design as in generate endpoint
    primary_color = RGBColor(0, 22, 57)
    secondary_color = RGBColor(59, 130, 246)
    text_color = RGBColor(71, 85, 105)
    
    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add background rectangle for visual appeal
        if idx == 0:  # Title slide gets special treatment
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(10), Inches(2.5)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = primary_color
            bg_shape.line.fill.background()
        else:
            # Subtle top accent bar for other slides
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(10), Inches(0.3)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = secondary_color
            accent_bar.line.fill.background()
        
        # Add title with better styling
        if slide_data.get("title"):
            title_y = Inches(0.6) if idx == 0 else Inches(0.5)
            title_height = Inches(1.2) if idx == 0 else Inches(0.8)
            title_box = slide.shapes.add_textbox(Inches(0.75), title_y, Inches(8.5), title_height)
            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44) if idx == 0 else Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255) if idx == 0 else primary_color
            title_para.alignment = PP_ALIGN.LEFT
            title_para.space_after = Pt(0)
        
        # Add content with better formatting
        if slide_data.get("content"):
            content_y = Inches(2.8) if idx == 0 else Inches(1.6)
            content_height = Inches(4.5) if idx == 0 else Inches(5.5)
            content_box = slide.shapes.add_textbox(Inches(0.75), content_y, Inches(8.5), content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.margin_left = Inches(0)
            content_frame.margin_right = Inches(0)
            content_frame.margin_top = Inches(0)
            content_frame.margin_bottom = Inches(0)
            
            # Parse content into paragraphs
            content_lines = slide_data["content"].split('\n')
            for line_idx, line in enumerate(content_lines):
                if line.strip():
                    if line_idx > 0:
                        p = content_frame.add_paragraph()
                    else:
                        p = content_frame.paragraphs[0]
                    
                    # Handle bullet points
                    if line.strip().startswith('•') or line.strip().startswith('-'):
                        p.text = line.strip().lstrip('•-').strip()
                        p.level = 0
                        p.font.size = Pt(20) if idx == 0 else Pt(18)
                        p.font.color.rgb = RGBColor(255, 255, 255) if idx == 0 else text_color
                        p.space_after = Pt(14)
                        p.space_before = Pt(0)
                    else:
                        p.text = line.strip()
                        p.level = 0
                        p.font.size = Pt(20) if idx == 0 else Pt(18)
                        p.font.color.rgb = RGBColor(255, 255, 255) if idx == 0 else text_color
                        p.space_after = Pt(14)
                        p.space_before = Pt(0)
        
        # Add slide number (except on title slide)
        if idx > 0:
            slide_num_shape = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
            slide_num_frame = slide_num_shape.text_frame
            slide_num_frame.text = f"{idx}"
            slide_num_para = slide_num_frame.paragraphs[0]
            slide_num_para.font.size = Pt(14)
            slide_num_para.font.color.rgb = RGBColor(148, 163, 184)
            slide_num_para.alignment = PP_ALIGN.RIGHT
    
    # Save to bytes
    deck_bytes = io.BytesIO()
    prs.save(deck_bytes)
    deck_bytes.seek(0)
    
    filename = f"{plan.get('name', 'pitch_deck')}.pptx"
    
    return Response(
        content=deck_bytes.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )

async def _generate_slides_content(
    plan: Dict,
    sections: List[Dict],
    financial_model: Optional[Dict],
    intake_data: Dict
) -> List[Dict]:
    """Generate slide content using AI"""
    
    try:
        business_name = intake_data.get("business_name", "Business")
        description = intake_data.get("business_description", "")
        value_prop = intake_data.get("unique_value_proposition", "")
        target_customers = intake_data.get("target_customers", "")
        revenue_model = intake_data.get("revenue_model", [])
        
        # Get key sections
        exec_summary = next((s for s in sections if s.get("section_type") == "executive_summary"), None)
        market_analysis = next((s for s in sections if s.get("section_type") == "market_analysis"), None)
        team_section = next((s for s in sections if s.get("section_type") == "team"), None)
        
        prompt = f"""Create a pitch deck for this business plan. Generate 8-10 slides in JSON format.

Business: {business_name}
Description: {description}
Value Proposition: {value_prop}
Target Customers: {target_customers}
Revenue Model: {', '.join(revenue_model) if revenue_model else 'Not specified'}

Executive Summary: {exec_summary.get('content', '')[:500] if exec_summary else 'N/A'}
Market Analysis: {market_analysis.get('content', '')[:500] if market_analysis else 'N/A'}
Team: {team_section.get('content', '')[:500] if team_section else 'N/A'}

Financial Highlights:
{_format_financial_highlights(financial_model) if financial_model else 'N/A'}

Generate slides in this JSON format:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "content": "Bullet points or short paragraphs for the slide content. Keep it concise (3-5 bullets max)."
    }}
  ]
}}

Include these slides:
1. Title slide (Business name and tagline)
2. Problem (What problem are you solving?)
3. Solution (Your product/service)
4. Market Opportunity (Market size and growth)
5. Business Model (How you make money)
6. Traction/Metrics (Key achievements or projections)
7. Team (Key team members)
8. Financials (Key financial highlights)
9. Ask/Use of Funds (What you're asking for and why)
10. Thank You/Contact

Return ONLY valid JSON, no other text."""
        
        chat = LlmChat(
            api_key=os.environ.get("OPENAI_API_KEY"),
            session_id=f"pitch_deck_{plan.get('_id')}",
            system_message="You are an expert pitch deck creator. Generate concise, compelling slide content in JSON format only."
        ).with_model("openai", "gpt-4o")
        
        response = await chat.send_message(UserMessage(text=prompt))
        
        # Parse JSON
        import json
        import re
        
        json_match = re.search(r'\{.*\}', response, re.DOTALL)
        if json_match:
            slides_data = json.loads(json_match.group())
            return slides_data.get("slides", [])
        else:
            # Fallback: create basic slides
            return _create_fallback_slides(intake_data, financial_model)
            
    except Exception as e:
        logger.error(f"Error generating slides: {e}")
        return _create_fallback_slides(intake_data, financial_model)

def _format_financial_highlights(financial_model: Dict) -> str:
    """Format financial model for prompt"""
    data = financial_model.get("data", {})
    pnl = data.get("pnl_monthly", [])
    
    highlights = []
    if pnl:
        first_year_revenue = sum([m.get("revenue", 0) for m in pnl[:12]])
        highlights.append(f"Year 1 Revenue: ${first_year_revenue:,.0f}")
    
    break_even = data.get("break_even", {})
    if break_even and break_even.get("months_to_break_even"):
        highlights.append(f"Break-even: {break_even['months_to_break_even']} months")
    
    return "\n".join(highlights) if highlights else "Financial projections available"

def _create_fallback_slides(intake_data: Dict, financial_model: Optional[Dict]) -> List[Dict]:
    """Create basic slides if AI generation fails"""
    business_name = intake_data.get("business_name", "Business")
    description = intake_data.get("business_description", "")
    value_prop = intake_data.get("unique_value_proposition", "")
    
    return [
        {
            "title": business_name,
            "content": "Your Business Pitch Deck"
        },
        {
            "title": "The Problem",
            "content": "• Identify the key problem your business solves\n• Explain why this problem matters\n• Show the pain points of your target customers"
        },
        {
            "title": "Our Solution",
            "content": f"{value_prop}\n\n{description}"
        },
        {
            "title": "Market Opportunity",
            "content": "• Market size and growth potential\n• Target customer segments\n• Market trends and opportunities"
        },
        {
            "title": "Business Model",
            "content": "• How you generate revenue\n• Pricing strategy\n• Revenue streams"
        },
        {
            "title": "Financial Highlights",
            "content": "• Key financial projections\n• Revenue growth\n• Break-even timeline"
        },
        {
            "title": "Our Team",
            "content": "• Key team members\n• Relevant experience\n• Why this team can execute"
        },
        {
            "title": "The Ask",
            "content": "• Funding amount requested\n• Use of funds\n• Expected milestones"
        },
        {
            "title": "Thank You",
            "content": f"{business_name}\nContact us to learn more"
        }
    ]

def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)
